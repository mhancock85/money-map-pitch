from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

def create_presentation():
    prs = Presentation()
    
    # 16:9 aspect ratio
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    # Colors
    bg_color = RGBColor(10, 31, 28)       # #0A1F1C
    primary_color = RGBColor(212, 242, 62)# #D4F23E
    text_white = RGBColor(255, 255, 255)
    text_grey = RGBColor(174, 189, 187)   # #aebdbb

    def set_slide_background(slide):
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = bg_color

    def add_title(slide, text, font_size=60):
        tbox = slide.shapes.add_textbox(Inches(1), Inches(0.5), Inches(11.33), Inches(1.5))
        tf = tbox.text_frame
        p = tf.paragraphs[0]
        p.text = text
        p.alignment = PP_ALIGN.LEFT
        p.font.name = 'Arial'
        p.font.size = Pt(font_size)
        p.font.color.rgb = primary_color
        p.font.bold = True

    def add_text(slide, text, top, left=1, width=11, font_size=24, color=text_white, bold=False):
        textbox = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(1))
        tf = textbox.text_frame
        p = tf.paragraphs[0]
        p.text = text
        p.font.name = 'Arial'
        p.font.size = Pt(font_size)
        p.font.color.rgb = color
        p.font.bold = bold
        return tf

    # ----------------------------------------------------------------
    # SLIDE 1: COVER
    # ----------------------------------------------------------------
    slide1 = prs.slides.add_slide(prs.slide_layouts[6]) # Blank
    set_slide_background(slide1)
    
    # Big Title
    tb = slide1.shapes.add_textbox(Inches(0), Inches(2.5), Inches(13.333), Inches(2))
    tf = tb.text_frame
    p = tf.paragraphs[0]
    p.text = "Money Map"
    p.font.name = 'Arial'
    p.font.size = Pt(96)
    p.font.color.rgb = text_white
    p.font.bold = True
    p.alignment = PP_ALIGN.CENTER

    # Subtitle
    tb2 = slide1.shapes.add_textbox(Inches(0), Inches(4), Inches(13.333), Inches(1))
    tf2 = tb2.text_frame
    p2 = tf2.paragraphs[0]
    p2.text = "The Financial Operating System for Global Citizens"
    p2.font.name = 'Arial'
    p2.font.size = Pt(32)
    p2.font.color.rgb = text_grey
    p2.alignment = PP_ALIGN.CENTER
    
    # Tag
    tb3 = slide1.shapes.add_textbox(Inches(5.6), Inches(5.5), Inches(2), Inches(0.8))
    tf3 = tb3.text_frame
    p3 = tf3.paragraphs[0]
    p3.text = "MVP PITCH"
    p3.font.size = Pt(14)
    p3.font.color.rgb = primary_color
    p3.alignment = PP_ALIGN.CENTER
    # (Cannot easily do rounded bg in pure python-pptx without complex shapes)

    # ----------------------------------------------------------------
    # SLIDE 2: THE PROBLEM
    # ----------------------------------------------------------------
    slide2 = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_background(slide2)
    add_title(slide2, "The Problem")
    
    add_text(slide2, "Financial management is broken for the international generation.", 1.5, color=text_grey)
    
    # 3 Columns
    col_width = 3.5
    gap = 0.5
    start_x = 1.2
    y_pos = 3.0
    
    # Col 1
    add_text(slide2, "Fragmented Finances", y_pos, start_x, col_width, 24, primary_color, True)
    add_text(slide2, "Bank accounts in UK, AU, and US. Credit cards everywhere. No single source of truth.", y_pos + 0.6, start_x, col_width, 18, text_grey)
    
    # Col 2
    add_text(slide2, "Manual Chaos", y_pos, start_x + col_width + gap, col_width, 24, primary_color, True)
    add_text(slide2, "Hours spent downloading PDFs and updating complex spreadsheets that break.", y_pos + 0.6, start_x + col_width + gap, col_width, 18, text_grey)

    # Col 3
    add_text(slide2, "Dumb Categorization", y_pos, start_x + (col_width + gap)*2, col_width, 24, primary_color, True)
    add_text(slide2, "Traditional apps don't know that 'Amazon' was a business expense, not a gift.", y_pos + 0.6, start_x + (col_width + gap)*2, col_width, 18, text_grey)


    # ----------------------------------------------------------------
    # SLIDE 3: THE SOLUTION
    # ----------------------------------------------------------------
    slide3 = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_background(slide3)
    add_title(slide3, "The Solution")
    
    add_text(slide3, "\"Set and Forget\" Intelligence", 1.8, font_size=32, color=primary_color, bold=True)
    add_text(slide3, "An intelligent layer that sits above your banks, unifying your financial life.", 2.5, color=text_grey, font_size=20)
    
    bullet_y = 3.5
    add_text(slide3, "1. Drop PDF Statements (Wise, Revolut, Banks)", bullet_y, font_size=24)
    add_text(slide3, "2. AI Extracts & Categorizes Contextually", bullet_y + 0.8, font_size=24)
    add_text(slide3, "3. Instant Premium Dashboard", bullet_y + 1.6, font_size=24)
    
    # Placeholder for visual
    viz_box = slide3.shapes.add_shape(1, Inches(8), Inches(2), Inches(4), Inches(4)) # msoShapeRectangle
    viz_box.fill.solid()
    viz_box.fill.fore_color.rgb = RGBColor(255, 255, 255)
    viz_box.fill.transparency = 0.9
    viz_box.line.color.rgb = primary_color
    
    tf = viz_box.text_frame
    p = tf.paragraphs[0]
    p.text = "Premium Dashboard Visual"
    p.font.color.rgb = width = primary_color


    # ----------------------------------------------------------------
    # SLIDE 4: BUSINESS MODEL
    # ----------------------------------------------------------------
    slide4 = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_background(slide4)
    add_title(slide4, "Business Model")
    add_text(slide4, "SaaS for the Modern Money Coach.", 1.5, color=text_grey)
    
    # Left Side
    add_text(slide4, "The \"Money Coach\" OS", 3, 1, 5, 28, primary_color, True)
    add_text(slide4, "Coaches currently spend 40% of their time chasing clients for spreadsheets. We automate data entry.", 3.6, 1, 5, 20, text_grey)
    add_text(slide4, "Subscription Model: £29 / mo / client", 5.5, 1, 5, 24, text_white, True)
    
    # Right Side
    add_text(slide4, "Why Wise & Burning Heroes?", 3, 7, 5, 24, text_white, True)
    add_text(slide4, "Wise Integration", 3.8, 7, 5, 20, primary_color, True)
    add_text(slide4, "Perfect use case for Wise API - multi-currency cashflow.", 4.2, 7, 5, 18, text_grey)
    
    add_text(slide4, "Scalability", 5.0, 7, 5, 20, primary_color, True)
    add_text(slide4, "Cloud-native (Modal + Supabase) architecture scales to zero.", 5.4, 7, 5, 18, text_grey)


    # ----------------------------------------------------------------
    # SLIDE 5: THE ASK
    # ----------------------------------------------------------------
    slide5 = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_background(slide5)
    
    # Title Centered
    tb = slide5.shapes.add_textbox(Inches(0), Inches(2), Inches(13.333), Inches(1.5))
    tf = tb.text_frame
    p = tf.paragraphs[0]
    p.text = "Join Us"
    p.font.size = Pt(72)
    p.font.bold = True
    p.font.color.rgb = text_white
    p.alignment = PP_ALIGN.CENTER
    
    add_text(slide5, "We are building the financial clarity layer for the next generation.", 3.5, 0, 13.333, 24, text_grey)
    # Hack to center the text above by creating a full width box, alignment default left in helper func but could be changed.
    # Actually the helper func makes text left aligned. Let's just manually fix alignment for this box if needed or assume user can edit.
    
    # 2 Cards for Ask
    box_y = 5.0
    box_w = 4.0
    box_h = 2.0
    gap = 1.0
    start_x = (13.333 - (box_w * 2 + gap)) / 2
    
    # Box 1
    add_text(slide5, "Seeking Mentorship", box_y, start_x, box_w, 20, text_grey)
    add_text(slide5, "Technical & Strategy", box_y + 0.5, start_x, box_w, 28, primary_color, True)
    
    # Box 2
    add_text(slide5, "Validation", box_y, start_x + box_w + gap, box_w, 20, text_grey)
    add_text(slide5, "Beta Users (Community)", box_y + 0.5, start_x + box_w + gap, box_w, 28, primary_color, True)


    # Output
    output_file = "Money_Map_Pitch_Deck.pptx"
    prs.save(output_file)
    print(f"Created {output_file}")

if __name__ == "__main__":
    create_presentation()
